"""
Gera a apresentacao PPTX do projeto:
"Sistema Multiagente de Elicitacao de Requisitos: Stakeholders como Personas de LLM".

Saida: apresentacao_empresa_junior.pptx  (abrir no Google Slides / PowerPoint / Keynote).
Design escuro, moderno, inspirado em keynotes Apple/Google.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- Paleta --------
BG      = RGBColor(0x0B, 0x0F, 0x1A)
BG_ALT  = RGBColor(0x11, 0x18, 0x28)
LINE    = RGBColor(0x22, 0x2C, 0x40)
TEXT    = RGBColor(0xF5, 0xF7, 0xFA)
MUTED   = RGBColor(0x9A, 0xA5, 0xB8)
ACCENT  = RGBColor(0x60, 0xA5, 0xFA)  # azul
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)  # verde
ACCENT3 = RGBColor(0xF4, 0x72, 0x6C)  # coral

FONT_T = "Inter"
FONT_B = "Inter"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=18, bold=False, italic=False,
        color=TEXT, font=FONT_B, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, ls=1.25):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
    return tb


def rich(s, x, y, w, h, paras, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, ls=1.3):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        if isinstance(para, str):
            para = [(para, {})]
        for text, opts in para:
            r = p.add_run(); r.text = text
            r.font.name = opts.get("font", FONT_B)
            r.font.size = Pt(opts.get("size", 16))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", TEXT)
    return tb


def rect(s, x, y, w, h, fill=BG_ALT, line=None, corner=None):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shape_kind, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if corner:
        sh.adjustments[0] = corner
    return sh


def hline(s, x1, y1, x2, y2, color=LINE, w=1.0):
    c = s.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c


def header(s, label, n, total):
    hline(s, Inches(0.6), Inches(0.55), Inches(12.73), Inches(0.55), color=LINE, w=0.75)
    txt(s, Inches(0.6), Inches(0.2), Inches(9), Inches(0.35),
        label.upper(), size=10, bold=True, color=ACCENT, font=FONT_T)
    txt(s, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.35),
        f"{n:02d} / {total:02d}", size=10, color=MUTED, font=FONT_T,
        align=PP_ALIGN.RIGHT)


def footer(s):
    txt(s, Inches(0.6), Inches(7.1), Inches(10), Inches(0.3),
        "Sistema Multiagente de Elicitacao de Requisitos  ·  Empresa Júnior · UFRJ",
        size=9, color=MUTED, font=FONT_T)


def placeholder(s, x, y, w, h, label):
    box = rect(s, x, y, w, h, fill=BG_ALT, line=ACCENT, corner=0.04)
    ln = box.line._get_or_add_ln()
    dash = etree.SubElement(ln, qn('a:prstDash'))
    dash.set('val', 'dash')
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT_T; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "\n[ inserir imagem editavel aqui ]"
    r2.font.name = FONT_B; r2.font.size = Pt(11); r2.font.italic = True
    r2.font.color.rgb = MUTED


TOTAL = 18

# ============================================================
# 1  ·  Capa
# ============================================================
s = slide()
rect(s, 0, 0, Inches(0.18), SH, fill=ACCENT)
txt(s, Inches(0.9), Inches(1.0), Inches(6), Inches(0.4),
    "UFRJ  ·  ENGENHARIA DE SOFTWARE  ·  2025",
    size=11, bold=True, color=ACCENT, font=FONT_T)
rich(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(3),
     [[("Sistema Multiagente de", {"size": 54, "bold": True, "font": FONT_T})],
      [("Elicitacao de Requisitos",
        {"size": 54, "bold": True, "font": FONT_T, "color": ACCENT})]],
     ls=1.05)
txt(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.6),
    "Stakeholders como Personas de LLM",
    size=24, color=MUTED, font=FONT_T)
hline(s, Inches(0.9), Inches(5.35), Inches(4.5), Inches(5.35), color=LINE)
rich(s, Inches(0.9), Inches(5.55), Inches(12), Inches(1.4),
     [[("Cassio Emanuel  ·  Giovanna Lavouras  ·  Rafael Maiani  ·  Rodrigo Nogueira  ·  Thiago Barcellos",
        {"size": 14, "color": TEXT, "font": FONT_T})],
      [("Estudo de caso: Empresa Júnior Consultoria — Empresa Junior da UFRJ",
        {"size": 13, "color": MUTED, "italic": True, "font": FONT_T})]],
     ls=1.4)

# ============================================================
# 2  ·  O Problema
# ============================================================
s = slide(); n = 2
header(s, "01  ·  Contexto", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.2),
    "A elicitacao de requisitos ainda e a fase mais critica —\ne tambem a mais fragil — do desenvolvimento de software.",
    size=32, bold=True, color=TEXT, font=FONT_T, ls=1.15)
rich(s, Inches(0.6), Inches(3.1), Inches(7.7), Inches(4),
     [[("Entrevistar stakeholders e caro, demorado e propenso a ruido. "
        "Informacoes se perdem, ambiguidades passam despercebidas e o esforco "
        "manual de analise qualitativa e enorme.", {"size": 17, "color": TEXT})],
      [(" ", {"size": 8})],
      [("Alem disso, a IA generativa aplicada a Engenharia de Software esta "
        "concentrada em codigo, testes e operacao ", {"size": 17, "color": TEXT}),
       ("— ha uma lacuna clara nas fases iniciais do ciclo, exatamente onde "
        "os erros custam mais caro.",
        {"size": 17, "color": ACCENT, "bold": True})]],
     ls=1.45)
rect(s, Inches(8.9), Inches(3.1), Inches(3.8), Inches(3.7), fill=BG_ALT, corner=0.06)
txt(s, Inches(9.15), Inches(3.25), Inches(3.4), Inches(0.4),
    "PONTOS DE DOR", size=10, bold=True, color=ACCENT, font=FONT_T)
rich(s, Inches(9.15), Inches(3.7), Inches(3.4), Inches(3),
     [[("Ambiguidade", {"size": 15, "bold": True})],
      [("Requisitos conflitantes entre stakeholders.",
        {"size": 12, "color": MUTED})],
      [(" ", {"size": 6})],
      [("Perda de informacao", {"size": 15, "bold": True})],
      [("Nuances das entrevistas nao chegam a especificacao.",
        {"size": 12, "color": MUTED})],
      [(" ", {"size": 6})],
      [("Esforco manual", {"size": 15, "bold": True})],
      [("Analise qualitativa consome dias de trabalho.",
        {"size": 12, "color": MUTED})]],
     ls=1.35)

# ============================================================
# 3  ·  Proposta
# ============================================================
s = slide(); n = 3
header(s, "02  ·  Proposta", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.4),
    "E se cada stakeholder fosse\numa persona interpretada por um agente de IA?",
    size=32, bold=True, color=TEXT, font=FONT_T, ls=1.15)
rich(s, Inches(0.6), Inches(3.3), Inches(12), Inches(3.5),
     [[("Nossa proposta: ", {"size": 18, "bold": True, "color": ACCENT}),
       ("um pipeline automatizado, baseado em IA generativa, que apoia "
        "a Engenharia de Requisitos por meio de um ", {"size": 18}),
       ("sistema multiagente orientado por personas.",
        {"size": 18, "bold": True})],
      [(" ", {"size": 8})],
      [("A partir de entrevistas estruturadas com stakeholders reais, "
        "o sistema gera personas fieis, instancia um agente para cada uma delas, "
        "produz requisitos individuais sob cada perspectiva e, finalmente, "
        "consolida tudo num Documento de Especificacao de Requisitos (SRS) "
        "com o apoio de ", {"size": 18, "color": TEXT}),
       ("Retrieval-Augmented Generation (RAG)",
        {"size": 18, "bold": True, "color": ACCENT2}),
       (" sobre a base de conhecimento da organizacao.",
        {"size": 18, "color": TEXT})]],
     ls=1.5)

# ============================================================
# 4  ·  Trabalhos Relacionados
# ============================================================
s = slide(); n = 4
header(s, "03  ·  Trabalhos Relacionados", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7),
    "O que ja existe — e o que ainda faltava.",
    size=28, bold=True, color=TEXT, font=FONT_T)
cards = [
    ("RECOVER  ·  Voria et al. (2025)",
     "Gera requisitos direto da conversa. Nao modela personas nem multiplas perspectivas.",
     ACCENT),
    ("Multiagentes  ·  Samil et al. (2024)",
     "PO, Dev, QA como agentes. Trabalha com cenarios abstratos, sem usuarios reais.",
     ACCENT),
    ("W6H  ·  Sultan & Miranskyy (2018)",
     "Padrao Who/What/Which/Where/How/Why/When para reduzir ambiguidade em entrevistas.",
     ACCENT2),
    ("Silicon Sampling  ·  Argyle (2023) / Anthis (2025)",
     "LLMs simulando humanos. Defende context-rich prompting com entrevistas reais.",
     ACCENT2),
]
xs = [Inches(0.6), Inches(6.9)]
ys = [Inches(2.0), Inches(4.55)]
for i, (t, b, col) in enumerate(cards):
    x, y = xs[i % 2], ys[i // 2]
    rect(s, x, y, Inches(6.0), Inches(2.3), fill=BG_ALT, corner=0.05)
    rect(s, x, y, Inches(0.08), Inches(2.3), fill=col)
    txt(s, x + Inches(0.35), y + Inches(0.25), Inches(5.5), Inches(0.5),
        t, size=15, bold=True, color=TEXT, font=FONT_T)
    txt(s, x + Inches(0.35), y + Inches(0.85), Inches(5.5), Inches(1.4),
        b, size=13, color=MUTED, font=FONT_B, ls=1.35)
txt(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
    "Nosso diferencial: unir personas fundamentadas em entrevistas reais + agentes autonomos + RAG num unico pipeline.",
    size=11, color=ACCENT, font=FONT_T, bold=True)

# ============================================================
# 5  ·  Arquitetura (visao geral)
# ============================================================
s = slide(); n = 5
header(s, "04  ·  Arquitetura", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
    "Visao geral: quatro etapas, um so pipeline.",
    size=28, bold=True, color=TEXT, font=FONT_T)
txt(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
    "Da entrevista bruta ao documento de requisitos consolidado.",
    size=15, color=MUTED, font=FONT_T)
steps = [
    ("01", "Coleta",       "Entrevistas estruturadas\ncom padrao W6H\n(Google Forms)"),
    ("02", "Personas",     "GPT-5 sintetiza personas\ndetalhadas por stakeholder\n(validacao Likert)"),
    ("03", "Simulador",    "Agente Gemini 2.5 Flash\ngera requisitos por persona"),
    ("04", "Consolidador", "Agente + RAG na base\nde conhecimento da Empresa Júnior\n=> SRS unificado"),
]
w = Inches(2.95); gap = Inches(0.25)
sx = Inches(0.6); sy = Inches(2.55)
for i, (num, t, b) in enumerate(steps):
    x = sx + (w + gap) * i
    rect(s, x, sy, w, Inches(3.4), fill=BG_ALT, corner=0.05)
    txt(s, x + Inches(0.3), sy + Inches(0.25), Inches(2.6), Inches(0.4),
        num, size=13, bold=True, color=ACCENT, font=FONT_T)
    txt(s, x + Inches(0.3), sy + Inches(0.75), Inches(2.6), Inches(0.6),
        t, size=22, bold=True, color=TEXT, font=FONT_T)
    txt(s, x + Inches(0.3), sy + Inches(1.55), Inches(2.6), Inches(1.8),
        b, size=13, color=MUTED, font=FONT_B, ls=1.4)
for i in range(3):
    x = sx + (w + gap) * (i + 1) - Inches(0.22)
    rect(s, x, sy + Inches(1.6), Inches(0.18), Inches(0.06), fill=ACCENT)
placeholder(s, Inches(0.6), Inches(6.15), Inches(12.15), Inches(0.65),
            "Figura 1 do relatorio — diagrama geral do pipeline (opcional)")

# ============================================================
# 6  ·  Coleta / W6H
# ============================================================
s = slide(); n = 6
header(s, "05  ·  Etapa 1  ·  Coleta", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Entrevistas estruturadas com o padrao W6H.",
    size=28, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.0), Inches(6.8), Inches(4.5),
     [[("Um formulario no Google Forms — construido sobre o padrao ",
        {"size": 16}),
       ("W6H", {"size": 16, "bold": True, "color": ACCENT}),
       (" (Sultan & Miranskyy, 2018) — captura, de forma padronizada, o perfil e as necessidades de cada participante.",
        {"size": 16})],
      [(" ", {"size": 8})],
      [("As respostas alimentam diretamente o restante do pipeline como contexto rico e estruturado, evitando ambiguidade desde a origem.",
        {"size": 16, "color": TEXT})],
      [(" ", {"size": 8})],
      [("Aplicado no estudo de caso com ", {"size": 16}),
       ("6 membros da Empresa Júnior ",
        {"size": 16, "bold": True, "color": ACCENT2}),
       ("em 5 cargos distintos.", {"size": 16})]],
     ls=1.4)
rect(s, Inches(7.8), Inches(2.0), Inches(4.95), Inches(4.5), fill=BG_ALT, corner=0.05)
txt(s, Inches(8.05), Inches(2.15), Inches(4.5), Inches(0.4),
    "PADRAO W6H", size=10, bold=True, color=ACCENT, font=FONT_T)
w6h = [("Who",   "Quem e o stakeholder?"),
       ("What",  "O que ele precisa?"),
       ("Which", "Quais alternativas existem?"),
       ("Where", "Onde o sistema atua?"),
       ("How",   "Como sera utilizado?"),
       ("Why",   "Por que isso importa?"),
       ("When",  "Quando e usado?")]
for i, (k, v) in enumerate(w6h):
    y = Inches(2.65 + i * 0.5)
    txt(s, Inches(8.05), y, Inches(1.1), Inches(0.4),
        k, size=14, bold=True, color=ACCENT2, font=FONT_T)
    txt(s, Inches(9.2), y, Inches(3.5), Inches(0.4),
        v, size=13, color=TEXT, font=FONT_B)

# ============================================================
# 7  ·  Personas (metodo)
# ============================================================
s = slide(); n = 7
header(s, "06  ·  Etapa 2  ·  Personas", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Cada entrevista vira uma persona rica e estruturada.",
    size=26, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(4.5),
     [[("O modelo ", {"size": 16}),
       ("GPT-5 ", {"size": 16, "bold": True, "color": ACCENT}),
       ("recebe a entrevista completa e, guiado por um prompt de engenharia rigorosa, produz uma persona com:",
        {"size": 16})],
      [(" ", {"size": 6})],
      [("Papel no negocio  ·  Perfil comportamental", {"size": 15, "color": TEXT})],
      [("Objetivos  ·  Necessidades explicitas e implicitas", {"size": 15, "color": TEXT})],
      [("Principais dificuldades vividas no dia a dia", {"size": 15, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Cada persona foi gerada ", {"size": 16}),
       ("5 vezes ", {"size": 16, "bold": True, "color": ACCENT2}),
       ("para avaliar consistencia do modelo. A melhor versao foi enviada ao proprio stakeholder para validacao (Likert de 4 pontos: clareza, completude, corretude).",
        {"size": 16})]],
     ls=1.4)
rect(s, Inches(7.4), Inches(2.0), Inches(5.35), Inches(4.5), fill=BG_ALT, corner=0.05)
txt(s, Inches(7.65), Inches(2.15), Inches(5), Inches(0.4),
    "PROMPT — GERACAO DE PERSONA", size=11, bold=True, color=ACCENT, font=FONT_T)
rich(s, Inches(7.65), Inches(2.6), Inches(4.9), Inches(3.8),
     [[("\"Gerar descricao detalhada de persona para o stakeholder responsavel pela entrevista...\"",
        {"size": 14, "italic": True, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Instrucoes principais:", {"size": 13, "bold": True, "color": ACCENT2})],
      [("Identificar necessidades explicitas e implicitas",
        {"size": 13, "color": MUTED})],
      [("Tracar perfil a partir dos padroes de resposta",
        {"size": 13, "color": MUTED})],
      [("Gerar descricao consistente com o stakeholder",
        {"size": 13, "color": MUTED})],
      [(" ", {"size": 6})],
      [("Formato de saida fixo: Persona, Papel, Perfil, Objetivos, Necessidades, Dificuldades.",
        {"size": 13, "color": TEXT})]],
     ls=1.35)

# ============================================================
# 8  ·  Exemplo de persona
# ============================================================
s = slide(); n = 8
header(s, "06  ·  Etapa 2  ·  Personas", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Exemplo real — Assessora de Gestao de Pessoas.",
    size=26, bold=True, color=TEXT, font=FONT_T)
txt(s, Inches(0.6), Inches(1.75), Inches(12), Inches(0.5),
    "Trecho da persona gerada automaticamente a partir da entrevista real.",
    size=14, color=MUTED, italic=True, font=FONT_T)
rect(s, Inches(0.6), Inches(2.4), Inches(6.05), Inches(4.4), fill=BG_ALT, corner=0.05)
txt(s, Inches(0.85), Inches(2.55), Inches(5.5), Inches(0.4),
    "PAPEL NO NEGOCIO", size=11, bold=True, color=ACCENT, font=FONT_T)
txt(s, Inches(0.85), Inches(3.0), Inches(5.6), Inches(3.6),
    "Responsavel por administrar todo o ciclo de vida dos membros da Empresa Júnior: "
    "ingresso de trainees, acompanhamento de efetivos, feedbacks, avaliacoes "
    "e treinamento. O sistema e sua principal ferramenta diaria para consolidar "
    "informacoes e reduzir tarefas operacionais repetitivas.",
    size=13, color=TEXT, font=FONT_B, ls=1.45)
rect(s, Inches(6.9), Inches(2.4), Inches(5.83), Inches(4.4), fill=BG_ALT, corner=0.05)
txt(s, Inches(7.15), Inches(2.55), Inches(5.5), Inches(0.4),
    "OBJETIVOS  &  PRINCIPAIS DORES", size=11, bold=True, color=ACCENT2, font=FONT_T)
rich(s, Inches(7.15), Inches(3.0), Inches(5.4), Inches(3.6),
     [[("Centralizar informacoes de membros e trainees em um so sistema",
        {"size": 13})],
      [("Eliminar multiplos logins (SSO)", {"size": 13})],
      [("Reduzir atividades manuais e repetitivas", {"size": 13})],
      [(" ", {"size": 6})],
      [("Dor principal:", {"size": 13, "bold": True, "color": ACCENT3})],
      [("informacoes fragmentadas entre Notion, Google Planilhas, Docs e Drive — multiplos logins, retrabalho e inconsistencias entre registros.",
        {"size": 13, "color": MUTED, "italic": True})]],
     ls=1.45)

# ============================================================
# 9  ·  Simulador
# ============================================================
s = slide(); n = 9
header(s, "07  ·  Etapa 3  ·  Simulador", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
    "Cada persona ganha um agente — e produz seus proprios requisitos.",
    size=24, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.1), Inches(6.5), Inches(4.7),
     [[("Implementacao em ", {"size": 16}),
       ("Python + Google AI Studio (Gemini 2.5 Flash)",
        {"size": 16, "bold": True, "color": ACCENT}), (".", {"size": 16})],
      [(" ", {"size": 6})],
      [("Cada execucao recebe uma persona especifica como contexto e um prompt de sistema que define o comportamento do agente: raciocinar como o stakeholder representado e gerar requisitos alinhados aos seus interesses.",
        {"size": 15, "color": TEXT})],
      [(" ", {"size": 8})],
      [("Rodamos ", {"size": 15}),
       ("5 vezes por persona ",
        {"size": 15, "bold": True, "color": ACCENT2}),
       ("para avaliar estabilidade e escolher a melhor geracao.", {"size": 15})],
      [(" ", {"size": 8})],
      [("Resultado desta etapa: um documento de requisitos por perspectiva — isolado, sem contaminacao entre stakeholders.",
        {"size": 15, "color": TEXT})]],
     ls=1.4)
rect(s, Inches(7.4), Inches(2.1), Inches(5.35), Inches(4.7), fill=BG_ALT, corner=0.05)
txt(s, Inches(7.65), Inches(2.25), Inches(5), Inches(0.4),
    "EXEMPLO  ·  GERENTE DE PROJETOS", size=11, bold=True,
    color=ACCENT, font=FONT_T)
rich(s, Inches(7.65), Inches(2.7), Inches(4.9), Inches(4.0),
     [[("RF01 · Dashboard Unificado", {"size": 14, "bold": True, "color": TEXT})],
      [("Visao consolidada de cada projeto com status, sprint, membros e milestones.",
        {"size": 12, "color": MUTED})],
      [(" ", {"size": 4})],
      [("RF02 · Gestao de Sprints e Backlog",
        {"size": 14, "bold": True, "color": TEXT})],
      [("Planejar sprints, priorizar backlog e acompanhar via Kanban.",
        {"size": 12, "color": MUTED})],
      [(" ", {"size": 4})],
      [("RF04 · Timesheet Integrado",
        {"size": 14, "bold": True, "color": TEXT})],
      [("Registrar tempo por tarefa (equivalente ao Clockify).",
        {"size": 12, "color": MUTED})],
      [(" ", {"size": 4})],
      [("RF05 · Alocacao e Gestao de Squads",
        {"size": 14, "bold": True, "color": TEXT})],
      [("Alocar membros a squads e papeis especificos por projeto.",
        {"size": 12, "color": MUTED})]],
     ls=1.3)

# ============================================================
# 10  ·  Consolidador + RAG
# ============================================================
s = slide(); n = 10
header(s, "08  ·  Etapa 4  ·  Consolidador", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
    "O especialista: consolida perspectivas com apoio de RAG.",
    size=26, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.1), Inches(7.0), Inches(4.7),
     [[("Um segundo agente ", {"size": 16}),
       ("(Gemini 2.5 Flash) ",
        {"size": 16, "bold": True, "color": ACCENT}),
       ("recebe os requisitos individuais de todas as personas e produz um unico Documento de Especificacao de Requisitos (SRS).",
        {"size": 16})],
      [(" ", {"size": 6})],
      [("Sua missao: ", {"size": 16, "bold": True, "color": ACCENT2}),
       ("identificar conflitos, resolver redundancias e alinhar ambiguidades — sem perder as nuances de cada perspectiva.",
        {"size": 16})],
      [(" ", {"size": 8})],
      [("Para nao 'inventar' regras da organizacao, o agente e enriquecido com ",
        {"size": 15}),
       ("Retrieval-Augmented Generation (RAG) ",
        {"size": 15, "bold": True, "color": ACCENT2}),
       ("sobre a base de conhecimento da Empresa Júnior: Codigo de Etica, Estatuto e documentos institucionais indexados em ChromaDB com embeddings Gemini.",
        {"size": 15})]],
     ls=1.4)
placeholder(s, Inches(7.9), Inches(2.1), Inches(4.85), Inches(4.7),
            "Figura 3 — Diagrama arquitetural (RAG + agentes)")

# ============================================================
# 11  ·  Engenharia de Prompt
# ============================================================
s = slide(); n = 11
header(s, "09  ·  Engenharia de Prompt", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Prompts como codigo: papel, instrucoes e formato de saida.",
    size=24, bold=True, color=TEXT, font=FONT_T)
txt(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
    "Todo prompt do pipeline segue o mesmo esqueleto — reduz alucinacao e garante saida estruturada.",
    size=14, color=MUTED, italic=True, font=FONT_T)
blocks = [
    ("PAPEL",
     "Definicao explicita de quem o modelo e (ex: 'voce e um Engenheiro de Requisitos senior...')."),
    ("CONTEXTO",
     "Entrevista, persona ou requisitos anteriores + trechos recuperados via RAG."),
    ("INSTRUCOES",
     "Passos objetivos: analisar, comparar, resolver conflitos, evitar genericos."),
    ("FORMATO DE SAIDA",
     "Estrutura fixa (Persona/Papel/Objetivos... ou RF/RNF numerados) para permitir automacao."),
]
for i, (t, b) in enumerate(blocks):
    x = Inches(0.6 + (i % 2) * 6.15)
    y = Inches(2.55 + (i // 2) * 2.2)
    rect(s, x, y, Inches(6.0), Inches(1.95), fill=BG_ALT, corner=0.05)
    txt(s, x + Inches(0.35), y + Inches(0.25), Inches(5.5), Inches(0.5),
        t, size=13, bold=True, color=ACCENT, font=FONT_T)
    txt(s, x + Inches(0.35), y + Inches(0.75), Inches(5.5), Inches(1.1),
        b, size=14, color=TEXT, font=FONT_B, ls=1.4)

# ============================================================
# 12  ·  Metodologia de avaliacao
# ============================================================
s = slide(); n = 12
header(s, "10  ·  Avaliacao", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
    "Como sabemos se funciona? Avaliacao hibrida em duas frentes.",
    size=24, bold=True, color=TEXT, font=FONT_T)
rect(s, Inches(0.6), Inches(2.1), Inches(6.05), Inches(4.9), fill=BG_ALT, corner=0.05)
txt(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(0.4),
    "A  ·  VALIDACAO DE PERSONAS", size=12, bold=True,
    color=ACCENT, font=FONT_T)
txt(s, Inches(0.85), Inches(2.85), Inches(5.5), Inches(0.6),
    "Feita pelos proprios stakeholders entrevistados.",
    size=15, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.85), Inches(3.55), Inches(5.5), Inches(3.3),
     [[("Escala Likert de 4 pontos avaliando tres criterios:",
        {"size": 14, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Clareza  ", {"size": 14, "bold": True, "color": ACCENT2}),
       ("— a persona e compreensivel?", {"size": 14, "color": TEXT})],
      [("Completude  ", {"size": 14, "bold": True, "color": ACCENT2}),
       ("— cobre o que o entrevistado disse?", {"size": 14, "color": TEXT})],
      [("Corretude  ", {"size": 14, "bold": True, "color": ACCENT2}),
       ("— fiel ao stakeholder real?", {"size": 14, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Nota minima 3 em cada criterio para aprovar. Caso contrario, regeramos a persona ou complementamos a entrevista.",
        {"size": 13, "color": MUTED})]],
     ls=1.4)
rect(s, Inches(6.9), Inches(2.1), Inches(5.83), Inches(4.9), fill=BG_ALT, corner=0.05)
txt(s, Inches(7.15), Inches(2.3), Inches(5.5), Inches(0.4),
    "B  ·  AVALIACAO DE REQUISITOS", size=12, bold=True,
    color=ACCENT, font=FONT_T)
txt(s, Inches(7.15), Inches(2.85), Inches(5.5), Inches(0.6),
    "Modelo INVENT — Engenharia de Requisitos.",
    size=15, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(7.15), Inches(3.55), Inches(5.5), Inches(3.3),
     [[("Inteligibilidade  ",
        {"size": 14, "bold": True, "color": ACCENT2}),
       ("— clareza do requisito", {"size": 14})],
      [("Necessidade  ",
        {"size": 14, "bold": True, "color": ACCENT2}),
       ("— resolve problema real?", {"size": 14})],
      [("Viabilidade  ",
        {"size": 14, "bold": True, "color": ACCENT2}),
       ("— e tecnica e praticamente implementavel?", {"size": 14})],
      [("Estabilidade  ",
        {"size": 14, "bold": True, "color": ACCENT2}),
       ("— resiste a mudancas de escopo?", {"size": 14})],
      [("Testabilidade  ",
        {"size": 14, "bold": True, "color": ACCENT2}),
       ("— tem criterios objetivos de validacao?", {"size": 14})]],
     ls=1.6)

# ============================================================
# 13  ·  Resultados personas
# ============================================================
s = slide(); n = 13
header(s, "11  ·  Resultados", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Personas — 6 stakeholders, 5 cargos, aprovacao total.",
    size=26, bold=True, color=TEXT, font=FONT_T)
rect(s, Inches(0.6), Inches(2.1), Inches(5.5), Inches(4.9), fill=BG_ALT, corner=0.06)
txt(s, Inches(0.85), Inches(2.4), Inches(5), Inches(0.5),
    "APROVACAO", size=12, bold=True, color=ACCENT, font=FONT_T)
txt(s, Inches(0.85), Inches(2.9), Inches(5), Inches(2.5),
    "6 / 6", size=140, bold=True, color=ACCENT2, font=FONT_T)
txt(s, Inches(0.85), Inches(5.6), Inches(5), Inches(1.3),
    "personas receberam nota maxima em todos os criterios (clareza, completude, corretude).",
    size=15, color=TEXT, font=FONT_B, ls=1.4)
rect(s, Inches(6.4), Inches(2.1), Inches(6.35), Inches(4.9), fill=BG_ALT, corner=0.06)
txt(s, Inches(6.65), Inches(2.4), Inches(5.8), Inches(0.4),
    "OBSERVACOES", size=12, bold=True, color=ACCENT, font=FONT_T)
rich(s, Inches(6.65), Inches(2.9), Inches(5.85), Inches(3.9),
     [[("Cada persona foi gerada 5 vezes ",
        {"size": 14, "bold": True, "color": TEXT}),
       ("para medir consistencia.", {"size": 14})],
      [(" ", {"size": 6})],
      [("As geracoes se mostraram estaveis e semanticamente equivalentes entre si — o modelo capturou de forma confiavel objetivos, dores e contexto de cada participante.",
        {"size": 14, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Cargos avaliados:",
        {"size": 13, "bold": True, "color": ACCENT2})],
      [("Assessora de Gestao de Pessoas · Assessora de Vendas + Dev Front-end · Lider Tecnico de Conhecimento · Lider Tecnico de Projetos · Gerente de Projetos.",
        {"size": 13, "color": MUTED})]],
     ls=1.4)

# ============================================================
# 14  ·  Resultados requisitos
# ============================================================
s = slide(); n = 14
header(s, "11  ·  Resultados", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "Requisitos — 49 gerados, 93,9% aprovados sem ajustes.",
    size=26, bold=True, color=TEXT, font=FONT_T)


def stat(s, x, big, big_color, label):
    rect(s, x, Inches(2.15), Inches(4.05), Inches(3.4), fill=BG_ALT, corner=0.06)
    txt(s, x + Inches(0.3), Inches(2.4), Inches(3.5), Inches(0.4),
        label, size=11, bold=True, color=ACCENT, font=FONT_T)
    txt(s, x + Inches(0.3), Inches(2.85), Inches(3.5), Inches(2.4),
        big, size=90, bold=True, color=big_color, font=FONT_T)


stat(s, Inches(0.6),  "49",    TEXT,    "REQUISITOS AVALIADOS")
stat(s, Inches(4.75), "93,9%", ACCENT2, "CORRETOS")
stat(s, Inches(8.9),  "6,1%",  ACCENT,  "CORRETOS COM AJUSTES")
txt(s, Inches(0.6), Inches(5.75), Inches(12), Inches(0.5),
    "Nenhum requisito foi considerado incorreto.",
    size=18, bold=True, color=TEXT, font=FONT_T)
txt(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.8),
    "Os ajustes sugeridos foram refinamentos contextuais: RF04 (busca com tags), RF14 (Gantt "
    "complementando o Kanban) e RF18 (adequacao ao processo real da Empresa Júnior).",
    size=14, color=MUTED, font=FONT_B, ls=1.4)

# ============================================================
# 15  ·  SRS final (exemplo)
# ============================================================
s = slide(); n = 15
header(s, "12  ·  Artefato Final", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
    "O SRS consolidado — trecho real do documento gerado.",
    size=24, bold=True, color=TEXT, font=FONT_T)
txt(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4),
    "36 requisitos funcionais em 6 modulos + 13 nao funcionais.  Arquivo: requisitos_consolidados_2.md",
    size=13, color=MUTED, italic=True, font=FONT_T)
rect(s, Inches(0.6), Inches(2.35), Inches(12.15), Inches(4.55), fill=BG_ALT, corner=0.03)
rect(s, Inches(0.6), Inches(2.35), Inches(0.08), Inches(4.55), fill=ACCENT2)
rich(s, Inches(0.9), Inches(2.55), Inches(11.7), Inches(4.3),
     [[("Modulo 1 — Plataforma, Seguranca e Colaboracao",
        {"size": 14, "bold": True, "color": ACCENT2})],
      [(" ", {"size": 4})],
      [("RF01 · Autenticacao Unificada (SSO)",
        {"size": 13, "bold": True, "color": TEXT}),
       ("  — acesso via credencial unica do Google Workspace.",
        {"size": 13, "color": MUTED})],
      [("RF02 · Controle de Acesso Baseado em Perfis (RBAC)",
        {"size": 13, "bold": True, "color": TEXT}),
       ("  — visualizacao/edicao por papel (Diretoria, GP, Vendas, Cliente...).",
        {"size": 13, "color": MUTED})],
      [("RF04 · Busca Global Indexada",
        {"size": 13, "bold": True, "color": TEXT}),
       ("  — localizar leads, projetos, atas e conteudo de anexos.",
        {"size": 13, "color": MUTED})],
      [(" ", {"size": 6})],
      [("Modulo 5 — Gestao de Pessoas & Conformidade Etica",
        {"size": 14, "bold": True, "color": ACCENT2})],
      [(" ", {"size": 4})],
      [("RF32 · Controle de Presenca e Sancoes Eticas",
        {"size": 13, "bold": True, "color": TEXT}),
       ("  — alerta automatico apos 3 faltas em RG (via Codigo de Etica, extraido pelo RAG).",
        {"size": 13, "color": MUTED})],
      [("RF33 · Direito ao Esquecimento (LGPD)",
        {"size": 13, "bold": True, "color": TEXT}),
       ("  — exclusao permanente de dados sensiveis mediante solicitacao.",
        {"size": 13, "color": MUTED})],
      [(" ", {"size": 6})],
      [("RNF01 · Conformidade com LGPD  ·  RNF04 · Busca em ate 2s  ·  RNF09 · Uptime 99,5%  ·  RNF11 · RTO 4h",
        {"size": 12, "italic": True, "color": ACCENT})]],
     ls=1.35)

# ============================================================
# 16  ·  Discussao
# ============================================================
s = slide(); n = 16
header(s, "13  ·  Discussao", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
    "O que esses resultados significam.",
    size=28, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.5),
     [[("O sistema consolidou corretamente perspectivas distintas — comerciais, tecnicas, operacionais e de gestao — sem gerar requisitos incompativeis com o dominio.",
        {"size": 17, "color": TEXT})],
      [(" ", {"size": 8})],
      [("As sugestoes de ajuste foram sempre ", {"size": 17}),
       ("refinamentos contextuais",
        {"size": 17, "bold": True, "color": ACCENT}),
       (", nunca correcoes estruturais. Isso sugere que a combinacao ", {"size": 17}),
       ("personas + agentes especializados + RAG ",
        {"size": 17, "bold": True, "color": ACCENT2}),
       ("preservou o contexto das entrevistas ao longo de todo o pipeline.",
        {"size": 17})],
      [(" ", {"size": 8})],
      [("Pontos de atencao observados: em raras execucoes, um mesmo requisito foi classificado como funcional em uma rodada e nao-funcional em outra; e algumas especificacoes foram alem do escopo razoavel de uma empresa junior (ex: exigir 99,5% de uptime).",
        {"size": 15, "color": MUTED})]],
     ls=1.45)

# ============================================================
# 17  ·  Conclusao & futuro
# ============================================================
s = slide(); n = 17
header(s, "14  ·  Conclusao", n, TOTAL); footer(s)
txt(s, Inches(0.6), Inches(0.9), Inches(12), Inches(1.0),
    "Conclusao e proximos passos.",
    size=28, bold=True, color=TEXT, font=FONT_T)
rich(s, Inches(0.6), Inches(2.1), Inches(6.4), Inches(4.7),
     [[("Contribuicao", {"size": 13, "bold": True, "color": ACCENT})],
      [(" ", {"size": 6})],
      [("Uma abordagem end-to-end de IA generativa para elicitacao de requisitos: personas fundamentadas em entrevistas reais + ecossistema multiagente + consolidacao com RAG.",
        {"size": 15, "color": TEXT})],
      [(" ", {"size": 8})],
      [("Validacao preliminar", {"size": 13, "bold": True, "color": ACCENT})],
      [(" ", {"size": 6})],
      [("6/6 personas aprovadas com nota maxima. 93,9% dos 49 requisitos considerados corretos. Nenhum requisito incorreto.",
        {"size": 15, "color": TEXT})]],
     ls=1.4)
rect(s, Inches(7.3), Inches(2.1), Inches(5.45), Inches(4.7), fill=BG_ALT, corner=0.05)
txt(s, Inches(7.55), Inches(2.3), Inches(5), Inches(0.4),
    "TRABALHOS FUTUROS", size=11, bold=True, color=ACCENT2, font=FONT_T)
rich(s, Inches(7.55), Inches(2.85), Inches(5), Inches(3.7),
     [[("Ampliar a avaliacao para outras organizacoes e mais stakeholders.",
        {"size": 14, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Comparar os requisitos gerados com especificacoes elaboradas por engenheiros de requisitos humanos.",
        {"size": 14, "color": TEXT})],
      [(" ", {"size": 6})],
      [("Explorar negociacao automatica entre agentes para resolver conflitos entre perspectivas durante a elicitacao.",
        {"size": 14, "color": TEXT})]],
     ls=1.45)

# ============================================================
# 18  ·  Obrigado / perguntas
# ============================================================
s = slide()
rect(s, 0, 0, Inches(0.18), SH, fill=ACCENT)
txt(s, Inches(0.9), Inches(1.2), Inches(6), Inches(0.4),
    "OBRIGADO", size=12, bold=True, color=ACCENT, font=FONT_T)
txt(s, Inches(0.9), Inches(2.0), Inches(12), Inches(2),
    "Perguntas?", size=96, bold=True, color=TEXT, font=FONT_T)
hline(s, Inches(0.9), Inches(4.3), Inches(4.5), Inches(4.3), color=LINE)
rich(s, Inches(0.9), Inches(4.5), Inches(12), Inches(2.5),
     [[("Sistema Multiagente de Elicitacao de Requisitos",
        {"size": 20, "bold": True, "color": TEXT, "font": FONT_T})],
      [("Stakeholders como Personas de LLM  ·  Estudo de caso: Empresa Júnior Consultoria",
        {"size": 15, "color": MUTED, "font": FONT_T})],
      [(" ", {"size": 8})],
      [("Cassio Emanuel  ·  Giovanna Lavouras  ·  Rafael Maiani  ·  Rodrigo Nogueira  ·  Thiago Barcellos",
        {"size": 13, "color": MUTED, "font": FONT_T})]],
     ls=1.4)


# ---------- salvar ----------
out = "apresentacao_empresa_junior.pptx"
prs.save(out)
print("OK:", out, "  slides:", len(prs.slides))
